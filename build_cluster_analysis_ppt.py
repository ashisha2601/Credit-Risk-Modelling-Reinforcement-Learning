#!/usr/bin/env python3
"""
Build a PowerPoint summarising synthetic-data clustering results.

Output:
  results/cluster_analysis.pptx

You can open this PPTX in PowerPoint/Keynote/Google Slides and export as PDF.
"""

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt


PROJECT_ROOT = Path(__file__).parent
RESULTS_DIR = PROJECT_ROOT / "results"


def add_title_slide(prs: Presentation) -> None:
    slide_layout = prs.slide_layouts[0]  # Title slide
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "Credit Risk Clustering on Synthetic Data"
    subtitle = slide.placeholders[1]
    subtitle.text = (
        "Kaggle-only • RBI priors-only • Hybrid (Kaggle + RBI)\n"
        "Auto-generated from ProjectSem5 results"
    )


def add_image_slide(
    prs: Presentation,
    title: str,
    image_path: Path,
    notes: str | None = None,
) -> None:
    if not image_path.exists():
        return

    slide_layout = prs.slide_layouts[5]  # Title Only
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title

    # Insert image filling most of the slide
    left = Inches(0.5)
    top = Inches(1.5)
    width = Inches(9)
    slide.shapes.add_picture(str(image_path), left, top, width=width)

    if notes:
        slide.notes_slide.notes_text_frame.text = notes


def build_presentation() -> Path:
    prs = Presentation()
    add_title_slide(prs)

    # ------------------------------------------------------------------
    # Overall (two-stage priors+CTGAN) 6-cluster view
    # ------------------------------------------------------------------
    add_image_slide(
        prs,
        "Two-Stage Priors+CTGAN: Cluster Sizes",
        RESULTS_DIR / "cluster_sizes.png",
        "Initial 6-cluster view on two-stage priors+CTGAN synthetic data.",
    )
    add_image_slide(
        prs,
        "Two-Stage Priors+CTGAN: Default Rate by Cluster",
        RESULTS_DIR / "cluster_default_rates.png",
        "Default rate per cluster for the two-stage priors+CTGAN dataset.",
    )
    add_image_slide(
        prs,
        "Two-Stage Priors+CTGAN: PCA Projection by Cluster",
        RESULTS_DIR / "cluster_pca_scatter.png",
        "PCA projection showing separation of the 6 clusters.",
    )

    # ------------------------------------------------------------------
    # Cross-dataset PD calibration
    # ------------------------------------------------------------------
    add_image_slide(
        prs,
        "Default Rate by Risk Segment Across Datasets",
        RESULTS_DIR / "cross_dataset_pd_by_risk_segment.png",
        "Comparison of LOW/MEDIUM/HIGH risk segments across Kaggle-only, "
        "RBI priors-only, and Hybrid datasets.",
    )

    # ------------------------------------------------------------------
    # Kaggle-only dataset
    # ------------------------------------------------------------------
    add_image_slide(
        prs,
        "Kaggle-only: Cluster Sizes",
        RESULTS_DIR / "kaggle_only_cluster_sizes.png",
    )
    add_image_slide(
        prs,
        "Kaggle-only: Default Rate by Cluster",
        RESULTS_DIR / "kaggle_only_cluster_default_rates.png",
        "All clusters have relatively high PD; Cluster 2 is highest.",
    )
    add_image_slide(
        prs,
        "Kaggle-only: Monthly Income by Cluster",
        RESULTS_DIR / "kaggle_only_box_AMT_INCOME_TOTAL.png",
    )
    add_image_slide(
        prs,
        "Kaggle-only: Loan Amount by Cluster",
        RESULTS_DIR / "kaggle_only_box_AMT_CREDIT.png",
    )
    add_image_slide(
        prs,
        "Kaggle-only: Monthly Annuity by Cluster",
        RESULTS_DIR / "kaggle_only_box_AMT_ANNUITY.png",
    )
    add_image_slide(
        prs,
        "Kaggle-only: Contract Type vs Cluster (Default Rate)",
        RESULTS_DIR / "kaggle_only_heatmap_cluster_vs_contract_type.png",
    )
    add_image_slide(
        prs,
        "Kaggle-only: Income vs EMI by Risk Segment",
        RESULTS_DIR / "kaggle_only_scatter_income_vs_emi.png",
    )

    # ------------------------------------------------------------------
    # RBI priors-only dataset
    # ------------------------------------------------------------------
    add_image_slide(
        prs,
        "RBI Priors-only: Cluster Sizes",
        RESULTS_DIR / "rbi_priors_only_cluster_sizes.png",
    )
    add_image_slide(
        prs,
        "RBI Priors-only: Default Rate by Cluster",
        RESULTS_DIR / "rbi_priors_only_cluster_default_rates.png",
        "Clean LOW/MEDIUM/HIGH risk ordering around 1–3% PD.",
    )
    add_image_slide(
        prs,
        "RBI Priors-only: PCA Projection by Cluster",
        RESULTS_DIR / "rbi_priors_only_cluster_pca_scatter.png",
        "Well-separated clusters generated directly from RBI priors.",
    )
    add_image_slide(
        prs,
        "RBI Priors-only: Monthly Income by Cluster",
        RESULTS_DIR / "rbi_priors_only_box_MONTHLY_INCOME.png",
    )
    add_image_slide(
        prs,
        "RBI Priors-only: Loan Amount by Cluster",
        RESULTS_DIR / "rbi_priors_only_box_LOAN_AMOUNT.png",
    )
    add_image_slide(
        prs,
        "RBI Priors-only: Credit Score by Cluster",
        RESULTS_DIR / "rbi_priors_only_box_CREDIT_SCORE.png",
    )
    add_image_slide(
        prs,
        "RBI Priors-only: Monthly Payment by Cluster",
        RESULTS_DIR / "rbi_priors_only_box_MONTHLY_PAYMENT.png",
    )
    add_image_slide(
        prs,
        "RBI Priors-only: Loan Type vs Cluster (Default Rate)",
        RESULTS_DIR / "rbi_priors_only_heatmap_cluster_vs_loan_type.png",
    )
    add_image_slide(
        prs,
        "RBI Priors-only: Bank Group vs Cluster (Default Rate)",
        RESULTS_DIR / "rbi_priors_only_heatmap_cluster_vs_bank_group.png",
    )
    add_image_slide(
        prs,
        "RBI Priors-only: Income vs EMI by Risk Segment",
        RESULTS_DIR / "rbi_priors_only_scatter_income_vs_emi.png",
    )

    # ------------------------------------------------------------------
    # Hybrid dataset
    # ------------------------------------------------------------------
    add_image_slide(
        prs,
        "Hybrid: Cluster Sizes",
        RESULTS_DIR / "integrated_hybrid_cluster_sizes.png",
    )
    add_image_slide(
        prs,
        "Hybrid: Default Rate by Cluster",
        RESULTS_DIR / "integrated_hybrid_cluster_default_rates.png",
        "Hybrid PDs around 2–3%, calibrated to RBI priors.",
    )
    add_image_slide(
        prs,
        "Hybrid: PCA Projection by Cluster",
        RESULTS_DIR / "integrated_hybrid_cluster_pca_scatter.png",
    )
    add_image_slide(
        prs,
        "Hybrid: Monthly Income by Cluster",
        RESULTS_DIR / "integrated_hybrid_box_MONTHLY_INCOME.png",
    )
    add_image_slide(
        prs,
        "Hybrid: Loan Amount by Cluster",
        RESULTS_DIR / "integrated_hybrid_box_LOAN_AMOUNT.png",
    )
    add_image_slide(
        prs,
        "Hybrid: Monthly Payment by Cluster",
        RESULTS_DIR / "integrated_hybrid_box_MONTHLY_PAYMENT.png",
    )
    add_image_slide(
        prs,
        "Hybrid: Loan Type vs Cluster (Default Rate)",
        RESULTS_DIR / "integrated_hybrid_heatmap_cluster_vs_loan_type.png",
    )
    add_image_slide(
        prs,
        "Hybrid: Bank Group vs Cluster (Default Rate)",
        RESULTS_DIR / "integrated_hybrid_heatmap_cluster_vs_bank_group.png",
    )
    add_image_slide(
        prs,
        "Hybrid: Income vs EMI by Risk Segment",
        RESULTS_DIR / "integrated_hybrid_scatter_income_vs_emi.png",
    )

    # Save PPTX
    RESULTS_DIR.mkdir(parents=True, exist_ok=True)
    output_path = RESULTS_DIR / "cluster_analysis.pptx"
    prs.save(output_path)
    return output_path


def main() -> None:
    output = build_presentation()
    print(f"Saved PowerPoint to: {output}")
    print("Open this PPTX and export as PDF from your presentation tool.")


if __name__ == "__main__":
    main()


